"""
Build The Abyss presentation as a .pptx file.
Run: python3 build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Palette ────────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x0B, 0x08)   # near-black
GOLD     = RGBColor(0xF5, 0xC8, 0x42)   # torch gold
RED      = RGBColor(0xC4, 0x2C, 0x2C)   # danger red
STONE    = RGBColor(0xA8, 0x98, 0x78)   # warm stone
DIM      = RGBColor(0x6B, 0x5E, 0x48)   # dimmer stone
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xCF, 0xC4, 0xAD)   # body text
CARD_BG  = RGBColor(0x1A, 0x16, 0x0F)   # card background
GREEN    = RGBColor(0x4A, 0x9A, 0x4A)   # solution green
CODE_BG  = RGBColor(0x16, 0x12, 0x0C)   # code block bg

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ────────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def bg(slide, color=BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1)):
    """Add a filled/outlined rectangle."""
    shp = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # rectangle
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.line.fill.background()
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        shp.line.width = line_w
    else:
        shp.line.fill.background()
    return shp

def txt(slide, text, x, y, w, h,
        size=24, bold=False, color=OFFWHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def heading(slide, text, y=0.25, size=36):
    """Gold heading with underline bar."""
    txt(slide, text, 0.55, y, 12.2, 0.7,
        size=size, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    box(slide, 0.55, y + 0.65, 12.2, 0.04, fill=DIM)

def bullet_tf(slide, items, x, y, w, h,
              size=18, color=OFFWHITE, spacing=1.15):
    """Multi-line text box from a list of (text, indent, bold, color) tuples."""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, str):
            text, indent, bold, col = item, 0, False, color
        else:
            text, indent, bold, col = item
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = indent
        from pptx.util import Pt as _Pt
        p.space_after = _Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = col if col else color
    return txb

def card(slide, x, y, w, h, title, body, title_color=GOLD, body_size=16):
    """Rounded-ish card with title + body."""
    box(slide, x, y, w, h, fill=CARD_BG, line=DIM, line_w=Pt(1.2))
    txt(slide, title, x+0.12, y+0.1,  w-0.2, 0.32,
        size=13, bold=True, color=title_color)
    txt(slide, body,  x+0.12, y+0.38, w-0.2, h-0.5,
        size=body_size, color=OFFWHITE, wrap=True)

def accent_bar(slide, x, y, h, color=RED):
    """Thin vertical accent bar."""
    box(slide, x, y, 0.06, h, fill=color)

def slide_num_label(slide, n, total=11):
    txt(slide, f"{n} / {total}", 12.6, 7.1, 0.7, 0.3,
        size=11, color=DIM, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

# Subtle top bar
box(s, 0, 0, 13.33, 0.08, fill=RGBColor(0x1A,0x14,0x08))
# Bottom bar
box(s, 0, 7.3, 13.33, 0.2, fill=RGBColor(0x1A,0x14,0x08))

# Decorative side lines
box(s, 0.4, 1.5, 0.04, 4.5, fill=DIM)
box(s, 12.9, 1.5, 0.04, 4.5, fill=DIM)

txt(s, "🕯  THE ABYSS  🕯", 0.5, 1.8, 12.3, 1.4,
    size=60, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

txt(s, "A Procedurally Generated Dungeon Crawler",
    0.5, 3.4, 12.3, 0.6,
    size=24, color=STONE, align=PP_ALIGN.CENTER)

box(s, 4.5, 4.15, 4.33, 0.04, fill=DIM)

txt(s, "Python Final Project  ·  2025–26",
    0.5, 4.3, 12.3, 0.4,
    size=16, color=DIM, align=PP_ALIGN.CENTER)

txt(s, "Group members:   _________________     _________________     _________________",
    0.5, 6.7, 12.3, 0.4,
    size=13, color=DIM, align=PP_ALIGN.CENTER)

slide_num_label(s, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What Is It?
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); heading(s, "What Is The Abyss?")

txt(s,
    "A top-down dungeon crawler where every floor is uniquely generated — "
    "no two playthroughs are ever the same. Navigate torch-lit halls, "
    "fight monsters, collect loot, and descend as deep as you dare.",
    0.55, 1.05, 12.2, 0.75, size=18, color=OFFWHITE)

# 4 feature cards
card(s, 0.55, 2.0,  2.9, 2.1, "🗺️  Procedural Maps",
     "Random rooms carved from solid rock, connected by hand-dug corridors. No two floors are identical.")
card(s, 3.6,  2.0,  2.9, 2.1, "👹  3 Enemy Types",
     "Goblins, Skeletons & Demons — each with unique speed, HP, damage, and sight range.")
card(s, 6.65, 2.0,  2.9, 2.1, "🕯️  Torch Lighting",
     "Radial light mask limits your field of view — darkness hides what lurks ahead.")
card(s, 9.7,  2.0,  3.1, 2.1, "📈  Infinite Floors",
     "Enemies scale in difficulty with each floor descended. How deep can you go?")

# Controls box
box(s, 0.55, 4.3, 12.2, 1.4, fill=CARD_BG, line=DIM)
txt(s, "Controls", 0.75, 4.38, 3, 0.35, size=14, bold=True, color=GOLD)
bullet_tf(s, [
    ("WASD / Arrow Keys  —  Move",             0, False, OFFWHITE),
    ("SPACE  —  Attack (circle radius around player)",  0, False, OFFWHITE),
    ("P  —  Pause / Resume          ESC  —  Quit",      0, False, OFFWHITE),
], 0.75, 4.72, 12.0, 0.85, size=16)

slide_num_label(s, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Why This Project?
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); heading(s, "Why This Project?")

rows = [
    ("The Challenge",  RED,
     "We wanted to build something genuinely impressive — not a simple script, but a real interactive experience."),
    ("The Idea",       GOLD,
     "Dungeon crawlers are a classic genre, but procedural generation makes every run different — a genuinely hard computer science problem."),
    ("The Ambition",   STONE,
     "Games are the perfect test of Python skills: graphics, algorithms, data structures, and real-time logic all in one project."),
    ("The Result",     GREEN,
     "~1,000 lines of documented Python — procedural generation, BFS pathfinding, ray-marching line-of-sight, and torch lighting, all from scratch."),
]
for i, (label, col, body) in enumerate(rows):
    y = 1.1 + i * 1.42
    accent_bar(s, 0.55, y, 1.15, color=col)
    txt(s, label, 0.75, y,       3.5, 0.38, size=15, bold=True, color=col)
    txt(s, body,  0.75, y+0.38, 12.0, 0.75, size=17, color=OFFWHITE)

slide_num_label(s, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); heading(s, "How It Works — Three Core Systems")

card(s, 0.55, 1.1, 3.9, 3.8,
     "🏗️  1 · Generation",
     "Random rooms placed on a 54×42 tile grid.\n\n"
     "L-shaped corridors carved between each adjacent pair — guaranteeing every room is reachable.\n\n"
     "Enemies and items scattered through rooms 1–N.\n\n"
     "Difficulty (enemy count & HP) scales with floor number.")

card(s, 4.7, 1.1, 3.9, 3.8,
     "🧠  2 · Enemy AI",
     "BFS shortest-path finds routes around walls.\n\n"
     "Ray-marching line-of-sight check: enemies only chase if they can actually see you.\n\n"
     "Paths cached and recalculated every 0.48 s — responsive but cheap.\n\n"
     "Three states: idle → chasing → attacking.")

card(s, 8.85, 1.1, 3.95, 3.8,
     "🎨  3 · Rendering",
     "Camera tracks player position smoothly.\n\n"
     "Tiles drawn with depth-based darkening — deeper floors are gloomier.\n\n"
     "Torch mask (pre-rendered radial gradient) multiplied over the full frame using BLEND_RGB_MULT.\n\n"
     "Minimap drawn in top-right corner every frame.")

txt(s, "All three systems run inside a single 60 FPS game loop — Events → Logic → Render → Display",
    0.55, 5.1, 12.2, 0.4, size=15, color=DIM, align=PP_ALIGN.CENTER)

slide_num_label(s, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Procedural Generation (code)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); heading(s, "Procedural Generation")

# Left: explanation steps
card(s, 0.55, 1.1, 4.2, 1.5, "Step 1 — Rooms",
     "Try 400 random placements. Accept if no overlap with existing rooms. Stop at 9–16 rooms.", body_size=15)
card(s, 0.55, 2.75, 4.2, 1.5, "Step 2 — Corridors",
     "Sort rooms left→right. Carve 2-tile-wide L-shaped paths between each adjacent pair.", body_size=15)
card(s, 0.55, 4.4, 4.2, 1.5, "Step 3 — Populate",
     "Place player in room[0], stairs in last room, enemies & loot in all rooms between.", body_size=15)

# Right: code block
box(s, 5.0, 1.05, 7.8, 4.95, fill=CODE_BG, line=DIM)
code = (
    "def generate_dungeon(floor_num):\n"
    '    tiles = [[WALL] * MAP_W for _ in range(MAP_H)]\n'
    "    rooms = []\n\n"
    "    for _ in range(400):          # try 400 placements\n"
    "        rw = random.randint(4, 12)\n"
    "        rh = random.randint(4, 8)\n"
    "        rx = random.randint(1, MAP_W - rw - 1)\n"
    "        ry = random.randint(1, MAP_H - rh - 1)\n\n"
    "        # skip if it overlaps an existing room\n"
    "        if any(rx < r['x']+r['w']+1 and\n"
    "               rx+rw > r['x']-1 and ...\n"
    "               for r in rooms): continue\n\n"
    "        rooms.append({'x':rx,'y':ry,'w':rw,'h':rh})\n"
    "        for ty in range(ry, ry + rh):     # carve room\n"
    "            for tx in range(rx, rx + rw):\n"
    "                tiles[ty][tx] = FLOOR\n\n"
    "    # connect rooms with L-shaped corridors\n"
    "    for i in range(len(rooms) - 1):\n"
    "        _carve(rooms[i], rooms[i+1])\n"
    "    return {'tiles': tiles, 'rooms': rooms, ...}"
)
txt(s, code, 5.1, 1.1, 7.6, 4.85,
    size=12, color=RGBColor(0xAB, 0xCB, 0xAB), wrap=False)

slide_num_label(s, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Enemy AI (code)
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); heading(s, "Enemy AI — BFS Pathfinding + Line of Sight")

txt(s,
    "Enemies are not telepathic — they can only chase you if they have line of sight. "
    "Once spotted, BFS finds the shortest walkable route.",
    0.55, 1.05, 12.2, 0.55, size=17, color=OFFWHITE)

# BFS code
box(s, 0.55, 1.75, 6.0, 3.5, fill=CODE_BG, line=DIM)
bfs_code = (
    "def bfs_path(tiles, start, goal):\n"
    "    q       = deque([(sx, sy, [])])\n"
    "    visited = {(sx, sy)}\n\n"
    "    while q:\n"
    "        x, y, path = q.popleft()\n"
    "        for dx, dy in ((1,0),(-1,0),\n"
    "                       (0,1),(0,-1)):\n"
    "            nx, ny = x+dx, y+dy\n"
    "            if tiles[ny][nx] == WALL:\n"
    "                continue\n"
    "            if nx==gx and ny==gy:\n"
    "                return path+[(nx,ny)]\n"
    "            q.append((nx,ny,path+\n"
    "                       [(nx,ny)]))\n"
    "    return []   # unreachable"
)
txt(s, bfs_code, 0.65, 1.83, 5.8, 3.35,
    size=13, color=RGBColor(0xAB, 0xCB, 0xAB), wrap=False)

# LOS code
box(s, 6.8, 1.75, 6.0, 3.5, fill=CODE_BG, line=DIM)
los_code = (
    "def has_los(tiles, ax, ay, bx, by):\n"
    "    dist  = math.hypot(bx-ax, by-ay)\n"
    "    steps = max(2, int(dist * 3))\n\n"
    "    for i in range(steps + 1):\n"
    "        t  = i / steps\n"
    "        tx = int(ax + (bx-ax) * t)\n"
    "        ty = int(ay + (by-ay) * t)\n\n"
    "        if tiles[ty][tx] == WALL:\n"
    "            return False  # blocked\n\n"
    "    return True           # clear"
)
txt(s, los_code, 6.9, 1.83, 5.8, 3.35,
    size=13, color=RGBColor(0xAB, 0xCB, 0xAB), wrap=False)

# Labels
txt(s, "Breadth-First Search",  0.55, 5.35, 6.0, 0.3, size=14, bold=True, color=GOLD)
txt(s, "Ray-Marching LOS",      6.8,  5.35, 6.0, 0.3, size=14, bold=True, color=GOLD)

txt(s, "Paths recalculated every 0.48 s per enemy — cheap enough to run on 20+ enemies simultaneously",
    0.55, 5.75, 12.2, 0.35, size=14, color=DIM, align=PP_ALIGN.CENTER)

slide_num_label(s, 6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Data Structures
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); heading(s, "Data Structures")

card(s, 0.55, 1.1, 5.9, 2.6,
     "🗂️  Tile Grid  —  list[list[int]]",
     "54 × 42 two-dimensional list. Each cell holds WALL=0, FLOOR=1, or STAIRS=2.\n\n"
     "O(1) random access by (x, y) coordinate — fast enough for collision detection every frame.\n\n"
     "Generated fresh for each floor, thrown away on descent.")

card(s, 6.65, 1.1, 6.1, 2.6,
     "👤  Player State  —  dict",
     "Single dictionary holds everything about the player:\n\n"
     "x, y  ·  hp, max_hp  ·  facing  ·  atk_t, atk_cd  ·  inv_t  ·  alive\n\n"
     "Passed into every draw function — one source of truth, easy to inspect and debug.")

card(s, 0.55, 3.9, 5.9, 2.7,
     "👹  Enemy List  —  list[dict]",
     "Each enemy is its own dict: type, hp, spd, dmg, sight, path, state, facing, alive.\n\n"
     "The list is looped each frame for AI updates, hit detection, and rendering.\n\n"
     "Dead enemies stay in the list with alive=False — avoids mid-loop deletion bugs.")

card(s, 6.65, 3.9, 6.1, 2.7,
     "✨  Particles  —  list[dict]",
     "Kill sparks and pickup effects. Each particle stores:\n\n"
     "x, y, vx, vy, life, col, size\n\n"
     "Updated and removed each frame when life ≤ 0.\n\n"
     "Demonstrates dynamic list manipulation at runtime.")

slide_num_label(s, 7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Challenges & Solutions
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); heading(s, "Challenges & Solutions")

challenges = [
    (
        "Black square around the player",
        "The torch light mask was a fixed-size surface blitted directly onto the game — visible dark corners formed outside the circle radius.",
        "Created a full-screen black surface, blitted the white radial gradient at the player's position, then multiplied it over the entire frame using pygame's BLEND_RGB_MULT flag."
    ),
    (
        "Wall collision misalignment",
        "Player's physics position was at the tile's top-left corner, but the sprite rendered half a tile ahead — causing visual overlap into walls.",
        "Spawned all entities at tile + 0.5 so their physics position IS their visual centre. Removed all manual offset adjustments from rendering."
    ),
    (
        "Enemy pathfinding performance",
        "Running BFS for every enemy every single frame would freeze the game as floors became busier.",
        "Cached each enemy's path. Only recomputed every 0.48 seconds — still feels responsive but cuts pathfinding cost by ~97%."
    ),
]

for i, (title, challenge, solution) in enumerate(challenges):
    y = 1.15 + i * 2.05
    # challenge
    accent_bar(s, 0.55, y, 0.75, color=RED)
    box(s, 0.62, y, 12.13, 0.75, fill=RGBColor(0x1E, 0x10, 0x10), line=RGBColor(0x5A,0x22,0x22))
    txt(s, f"✗  {title}", 0.75, y+0.05, 11.9, 0.3, size=15, bold=True, color=RED)
    txt(s, challenge, 0.75, y+0.35, 11.9, 0.38, size=14, color=OFFWHITE)
    # solution
    accent_bar(s, 0.55, y+0.78, 0.72, color=GREEN)
    box(s, 0.62, y+0.78, 12.13, 0.72, fill=RGBColor(0x0F, 0x1A, 0x0F), line=RGBColor(0x2A,0x5A,0x2A))
    txt(s, "✓  Fixed:", 0.75, y+0.82, 2.0, 0.28, size=14, bold=True, color=GREEN)
    txt(s, solution,   0.75, y+0.82, 11.9, 0.62, size=13, color=OFFWHITE)

slide_num_label(s, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Error Handling & Code Quality
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); heading(s, "Error Handling & Code Quality")

card(s, 0.55, 1.1, 5.9, 2.5,
     "⚠️  try / except — Font Loading",
     "System fonts vary across computers. If Georgia or Arial aren't found, "
     "the game falls back to pygame's built-in font automatically:\n\n"
     "try:\n"
     "    font = pygame.font.SysFont('Georgia', 68)\n"
     "except Exception:\n"
     "    font = pygame.font.Font(None, 74)\n\n"
     "The game never crashes due to missing fonts.", body_size=14)

card(s, 6.65, 1.1, 6.1, 2.5,
     "🛡️  Input Validation — Collision",
     "Before moving the player, all four corners of the bounding box are "
     "checked against map bounds AND wall tiles:\n\n"
     "def walkable(nx, ny):\n"
     "    for ox, oy in corners:\n"
     "        tx, ty = int(nx+ox), int(ny+oy)\n"
     "        if not (0 <= tx < MAP_W ...): return False\n"
     "        if tiles[ty][tx] == WALL: return False\n"
     "    return True", body_size=14)

card(s, 0.55, 3.8, 5.9, 2.5,
     "📝  Docstrings on Every Function",
     "Every function documents its purpose, parameters, and return value:\n\n"
     "def bfs_path(tiles, start, goal, max_dist=28):\n"
     '    """\n'
     "    Find the shortest path using BFS.\n"
     "    Parameters:\n"
     "        tiles (list[list[int]]): Tile grid.\n"
     "        start (tuple): Start coords.\n"
     '    Returns: list[tuple] or []\n'
     '    """', body_size=14)

card(s, 6.65, 3.8, 6.1, 2.5,
     "🔒  Dungeon Generation Fallback",
     "If random room placement stalls (< 5 rooms placed after 400 attempts), "
     "a hard-coded fallback layout is used automatically — the game never starts "
     "with a broken or unplayable map.\n\n"
     "Ensures the game is always in a valid, playable state regardless of "
     "random seed outcomes.", body_size=14)

slide_num_label(s, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Next Steps
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s); heading(s, "Next Steps & Future Improvements")

next_items = [
    ("01", "Weapons & Inventory",
     "Different sword types, bows, consumables — let the player build a loadout as they descend."),
    ("02", "Boss Enemies",
     "One oversized enemy per 5th floor with custom multi-phase attack patterns."),
    ("03", "Sound & Music",
     "Sword swings, enemy growls, torch crackle — audio makes the atmosphere complete."),
    ("04", "Persistent High Scores",
     "Save best floor and gold to a JSON file so records survive between sessions."),
    ("05", "Traps & Secret Rooms",
     "Hidden rooms, spike traps, locked doors — reward thorough exploration."),
    ("06", "vs. Commercial Games",
     "Similar concept to Enter the Gungeon / Binding of Isaac — ours is simpler but 100% hand-coded."),
]

cols = [(0.55, 1.1), (4.65, 1.1), (8.75, 1.1),
        (0.55, 3.8), (4.65, 3.8), (8.75, 3.8)]

for (x, y), (num, title, body) in zip(cols, next_items):
    box(s, x, y, 3.85, 2.5, fill=CARD_BG, line=DIM)
    txt(s, num, x+0.12, y+0.1, 0.7, 0.5, size=28, bold=True, color=GOLD)
    txt(s, title, x+0.68, y+0.12, 3.1, 0.38, size=15, bold=True, color=STONE)
    txt(s, body,  x+0.12, y+0.6,  3.6, 1.75, size=14, color=OFFWHITE)

slide_num_label(s, 10)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Q&A
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)

box(s, 0, 0,    13.33, 0.08, fill=RGBColor(0x1A,0x14,0x08))
box(s, 0, 7.3,  13.33, 0.2,  fill=RGBColor(0x1A,0x14,0x08))
box(s, 0.4, 1.2, 0.04, 5.1, fill=DIM)
box(s, 12.9, 1.2, 0.04, 5.1, fill=DIM)

txt(s, "🕯   Questions?   🕯", 0.5, 1.5, 12.3, 1.3,
    size=56, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

txt(s, "We built every line from scratch in Python using pygame",
    0.5, 3.0, 12.3, 0.5, size=20, color=STONE, align=PP_ALIGN.CENTER)

box(s, 4.2, 3.65, 4.9, 0.04, fill=DIM)

stat_items = [
    ("Language", "Python 3"),
    ("Library",  "pygame-ce"),
    ("Algorithm","BFS + Ray-march"),
    ("Lines",    "~1,000"),
]
for i, (label, val) in enumerate(stat_items):
    x = 0.9 + i * 2.9
    box(s, x, 3.85, 2.55, 1.4, fill=CARD_BG, line=DIM)
    txt(s, label, x+0.12, 3.92, 2.35, 0.35, size=12, bold=True, color=STONE, align=PP_ALIGN.CENTER)
    txt(s, val,   x+0.12, 4.3,  2.35, 0.7,  size=20, bold=True, color=GOLD,  align=PP_ALIGN.CENTER)

txt(s, "github.com/ethanblairs2006-stack/the-abyss",
    0.5, 6.6, 12.3, 0.4, size=14, color=DIM, align=PP_ALIGN.CENTER)

slide_num_label(s, 11)


# ── Save ───────────────────────────────────────────────────────────────────────
out = "The Abyss - Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
